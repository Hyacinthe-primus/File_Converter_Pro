"""Tests for the multi-engine fallback chains in the converter dispatch."""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from converter.converters import AdvancedConverterEngine

RTF_SAMPLE = (
    r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Helvetica;}}"
    r"\f0\fs24 Hello {\b bold} and {\i italic}\par"
    r"Second {\fs48 large} text\par}"
)


def _engine():
    return AdvancedConverterEngine()


def _assert_pdf(path):
    assert path.exists(), f"expected PDF at {path}"
    assert path.stat().st_size > 0
    with path.open("rb") as f:
        assert f.read(4) == b"%PDF"


def test_rtf_to_pdf_dispatch_chain():
    spec, ext = _engine()._DISPATCH["rtf_to_pdf"]
    assert ext == "pdf"
    assert spec == ("_rtf_to_pdf_native", "_rtf_to_pdf")


def test_xlsx_to_pdf_dispatch_chain():
    spec, ext = _engine()._DISPATCH["xlsx_to_pdf"]
    assert ext == "pdf"
    assert spec == ("_xlsx_to_pdf_com", "_xlsx_to_pdf")


def test_pptx_to_pdf_dispatch_chain():
    spec, ext = _engine()._DISPATCH["pptx_to_pdf"]
    assert ext == "pdf"
    assert spec == ("_pptx_to_pdf_com", "_pptx_to_pdf_native", "_pptx_to_pdf")


def test_single_engine_types_stay_strings():
    engine = _engine()
    for key in ("txt_to_pdf", "image_to_png", "audio_to_mp3"):
        spec, _ext = engine._DISPATCH[key]
        assert isinstance(spec, str)


def test_all_dispatch_engines_resolve():
    engine = _engine()
    for spec, _ext in engine._DISPATCH.values():
        names = spec if isinstance(spec, tuple) else (spec,)
        for name in names:
            assert hasattr(engine, name), f"dispatch engine {name!r} not found"


def test_convert_unknown_type():
    result = _engine().convert("nope_to_nope", "/a.in", "/out")
    assert result.success is False
    assert "Unknown type" in result.error


def test_rtf_to_pdf_uses_native_renderer(tmp_path, monkeypatch):
    src = tmp_path / "doc.rtf"
    src.write_text(RTF_SAMPLE, encoding="utf-8")
    engine = _engine()
    calls = []
    real = engine._rtf_to_pdf_native

    def _spy(s, d):
        calls.append("native")
        return real(s, d)

    monkeypatch.setattr(engine, "_rtf_to_pdf_native", _spy)
    monkeypatch.setattr(engine, "_rtf_to_pdf", lambda s, d: calls.append("internal") or True)
    result = engine.convert("rtf_to_pdf", str(src), str(tmp_path))
    assert result.success is True
    assert calls == ["native"]
    _assert_pdf(tmp_path / "doc.pdf")


def test_convert_falls_back_when_primary_returns_false(tmp_path, monkeypatch):
    src = tmp_path / "doc.rtf"
    src.write_text(RTF_SAMPLE, encoding="utf-8")
    engine = _engine()
    calls = []
    monkeypatch.setattr(engine, "_rtf_to_pdf_native", lambda s, d: calls.append("native") or False)
    result = engine.convert("rtf_to_pdf", str(src), str(tmp_path))
    assert result.success is True
    assert calls == ["native"]
    _assert_pdf(tmp_path / "doc.pdf")


def test_convert_falls_back_when_primary_raises(tmp_path, monkeypatch):
    src = tmp_path / "doc.rtf"
    src.write_text(RTF_SAMPLE, encoding="utf-8")
    engine = _engine()

    def _boom(s, d):
        raise RuntimeError("Office COM not available")

    monkeypatch.setattr(engine, "_rtf_to_pdf_native", _boom)
    result = engine.convert("rtf_to_pdf", str(src), str(tmp_path))
    assert result.success is True
    _assert_pdf(tmp_path / "doc.pdf")


def test_convert_fails_when_all_engines_fail(tmp_path, monkeypatch):
    src = tmp_path / "doc.rtf"
    src.write_text(RTF_SAMPLE, encoding="utf-8")
    engine = _engine()
    monkeypatch.setattr(engine, "_rtf_to_pdf_native", lambda s, d: False)
    monkeypatch.setattr(engine, "_rtf_to_pdf", lambda s, d: False)
    result = engine.convert("rtf_to_pdf", str(src), str(tmp_path))
    assert result.success is False
    assert "_rtf_to_pdf_native" in result.error
    assert "_rtf_to_pdf" in result.error


def test_convert_reports_failure_when_no_output_written(tmp_path, monkeypatch):
    src = tmp_path / "note.txt"
    src.write_text("hello", encoding="utf-8")
    engine = _engine()
    monkeypatch.setattr(engine, "_txt_to_pdf", lambda s, d: True)
    result = engine.convert("txt_to_pdf", str(src), str(tmp_path))
    assert result.success is False
    assert "produced no output" in result.error


def test_xlsx_to_pdf_falls_back_to_openpyxl(tmp_path, monkeypatch):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Age"])
    ws.append(["Alice", 30])
    ws.append(["Bob", 25])
    src = tmp_path / "data.xlsx"
    wb.save(src)

    engine = _engine()
    calls = []
    monkeypatch.setattr(engine, "_xlsx_to_pdf_com", lambda s, d: calls.append("com") or False)
    result = engine.convert("xlsx_to_pdf", str(src), str(tmp_path))
    assert result.success is True
    assert calls == ["com"]
    _assert_pdf(tmp_path / "data.pdf")


def test_pptx_to_pdf_falls_back_to_native(tmp_path, monkeypatch):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX"
    slide.placeholders[1].text = "Bullet one\nBullet two"
    src = tmp_path / "deck.pptx"
    prs.save(src)

    engine = _engine()
    calls = []
    real = engine._pptx_to_pdf_native

    def _spy(s, d):
        calls.append("native")
        return real(s, d)

    monkeypatch.setattr(engine, "_pptx_to_pdf_com", lambda s, d: calls.append("com") or False)
    monkeypatch.setattr(engine, "_pptx_to_pdf_native", _spy)
    result = engine.convert("pptx_to_pdf", str(src), str(tmp_path))
    assert result.success is True
    assert calls == ["com", "native"]
    _assert_pdf(tmp_path / "deck.pdf")


def test_office_com_wrappers_pass_app_name():
    engine = _engine()
    calls = []
    engine._office_to_pdf_com = lambda src, dst, app: calls.append((src, dst, app)) or True
    assert engine._xlsx_to_pdf_com("a.xlsx", "b.pdf") is True
    assert engine._pptx_to_pdf_com("c.pptx", "d.pdf") is True
    assert calls == [
        ("a.xlsx", "b.pdf", "Excel.Application"),
        ("c.pptx", "d.pdf", "PowerPoint.Application"),
    ]


def test_convert_reports_engine_used_primary(tmp_path, monkeypatch):
    src = tmp_path / "doc.rtf"
    src.write_text(RTF_SAMPLE, encoding="utf-8")
    engine = _engine()
    result = engine.convert("rtf_to_pdf", str(src), str(tmp_path))
    assert result.success is True
    assert result.engine_used == "_rtf_to_pdf_native"
    assert result.degraded is False


def test_convert_marks_degraded_on_fallback(tmp_path, monkeypatch):
    src = tmp_path / "doc.rtf"
    src.write_text(RTF_SAMPLE, encoding="utf-8")
    engine = _engine()
    monkeypatch.setattr(engine, "_rtf_to_pdf_native", lambda s, d: False)
    result = engine.convert("rtf_to_pdf", str(src), str(tmp_path))
    assert result.success is True
    assert result.engine_used == "_rtf_to_pdf"
    assert result.degraded is True


def test_convert_failure_has_no_engine(tmp_path, monkeypatch):
    src = tmp_path / "doc.rtf"
    src.write_text(RTF_SAMPLE, encoding="utf-8")
    engine = _engine()
    monkeypatch.setattr(engine, "_rtf_to_pdf_native", lambda s, d: False)
    monkeypatch.setattr(engine, "_rtf_to_pdf", lambda s, d: False)
    result = engine.convert("rtf_to_pdf", str(src), str(tmp_path))
    assert result.success is False
    assert result.engine_used == ""
    assert result.degraded is False


def test_conversion_result_repr_mentions_engine():
    from converter.helpers import ConversionResult

    r = ConversionResult(True, "/a.rtf", "/a.pdf", engine_used="_rtf_to_pdf_native")
    assert "_rtf_to_pdf_native" in repr(r)
