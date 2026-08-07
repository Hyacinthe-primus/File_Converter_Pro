"""OptimizationMixin — Office file optimization methods.

PDF optimization strategy (see optimize_pdf_file): tries Ghostscript first
(rebuilds the PDF from rendered content — best general compression, also
recovers files with corrupted xref tables), falls back to PyMuPDF
(garbage-collection + font subsetting, works fully offline), then pikepdf
(QPDF-based repair, also fully offline) if PyMuPDF's save doesn't validate.
A file is never handed back larger than the original — see the final size
guard in optimize_pdf_file.
"""

import os
from pathlib import Path
from datetime import datetime

class OptimizationMixin:
    """Mixin: office file optimization methods for FileConverterApp."""

    def optimize_office_files(self, office_files, optimization_type, quality_level, remove_metadata, compress_images, keep_backup):
        if hasattr(self, 'active_templates') and 'office_optimization' in self.active_templates:
            del self.active_templates['office_optimization']
        """Optimize office and image files"""
        output_dir = self.get_output_directory()
        if not output_dir:
            return

        IMAGE_EXTS  = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp', '.gif'}
        EXCEL_EXTS  = {'.xlsx', '.xls'}
        AUDIO_EXTS  = {'.mp3', '.wav', '.aac', '.flac', '.ogg'}
        VIDEO_EXTS  = {'.mp4', '.avi', '.mkv', '.mov', '.webm'}
        WEB_EXTS    = {'.json', '.html', '.htm'}
        EPUB_EXTS   = {'.epub'}

        self.show_progress(True, self.translate_text("len_off").format(len(office_files)))

        success_count = 0
        total_original_size = 0
        total_compressed_size = 0
        start_time = datetime.now()

        for i, file_path in enumerate(office_files):
            try:
                file_ext = Path(file_path).suffix.lower()
                original_size = os.path.getsize(file_path)
                total_original_size += original_size

                if keep_backup:
                    output_file = os.path.join(output_dir, f"optimized_{Path(file_path).name}")
                else:
                    output_file = file_path

                operation_start = datetime.now()

                if quality_level == 0:
                    compression_level = "high"
                elif quality_level == 1:
                    compression_level = "normal"
                elif quality_level == 2:
                    compression_level = "very_reduced"

                if file_ext == '.pdf':
                    success = self.optimize_pdf_file(file_path, output_file, compression_level, remove_metadata, compress_images)
                elif file_ext in ['.docx', '.doc']:
                    success = self.optimize_word_file(file_path, output_file, compression_level, remove_metadata, compress_images)
                elif file_ext in ['.pptx', '.ppt']:
                    success = self.optimize_powerpoint_file(file_path, output_file, compression_level, remove_metadata, compress_images)
                elif file_ext in EXCEL_EXTS:
                    success = self.optimize_excel_file(file_path, output_file, compression_level, remove_metadata)
                elif file_ext in IMAGE_EXTS:
                    success = self.optimize_image_file(file_path, output_file, quality_level)
                elif file_ext in AUDIO_EXTS:
                    success = self.optimize_av_file(file_path, output_file, quality_level, 'audio')
                elif file_ext in VIDEO_EXTS:
                    success = self.optimize_av_file(file_path, output_file, quality_level, 'video')
                elif file_ext in WEB_EXTS:
                    success = self.optimize_web_file(file_path, output_file, file_ext)
                elif file_ext in EPUB_EXTS:
                    success = self.optimize_epub_file(file_path, output_file, compress_images, quality_level)
                else:
                    success = False

                if success:
                    compressed_size = os.path.getsize(output_file) if os.path.exists(output_file) else original_size
                    total_compressed_size += compressed_size
                    operation_time = (datetime.now() - operation_start).total_seconds()
                    self.db_manager.add_conversion_record(
                        source_file=file_path, source_format=file_ext.upper().replace('.', ''),
                        target_file=output_file, target_format=file_ext.upper().replace('.', ''),
                        operation_type="office_optimization", file_size=original_size,
                        conversion_time=operation_time, success=True,
                        notes=f"Type: {optimization_type}, Quality: {quality_level}")
                    self.achievement_system.record_conversion("office_optimization", original_size, True)
                    success_count += 1
                self.progress_bar.setValue(int((i + 1) / len(office_files) * 100))
            except Exception as e:
                self.db_manager.add_conversion_record(
                    source_file=file_path, source_format=Path(file_path).suffix.upper().replace('.', ''),
                    target_file="", target_format="", operation_type="office_optimization",
                    file_size=0, conversion_time=0, success=False, notes=f"Error: {str(e)}")
                print(f"Optimization error {file_path}: {e}")

        total_time = (datetime.now() - start_time).total_seconds()
        self.show_progress(False)

        if total_original_size > 0:
            compression_rate = ((total_original_size - total_compressed_size) / total_original_size * 100)
            savings_mb = (total_original_size - total_compressed_size) / (1024 * 1024)
            message = self.translate_text("msg_1").format(success_count, len(office_files), f"{total_time:.1f}")
            message += self.translate_text("msg_2").format(f"{savings_mb:.2f}", f"{compression_rate:.1f}")
            message += self.translate_text("msg_3").format(output_dir)
        else:
            message = self.translate_text("msg_4").format(success_count, len(office_files))
        from PySide6.QtWidgets import QMessageBox
        QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))

    def optimize_pdf_file(self, pdf_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            original_size = os.path.getsize(pdf_path)

            # Ghostscript first: rebuilds the PDF from rendered content rather
            # than patching the existing object table, so it handles both
            # normal compression AND recovers files with corrupted xrefs in
            # one pass - no need to detect corruption and branch separately.
            print(f"[optimize_pdf] '{pdf_path}': trying Ghostscript first")
            if self._ghostscript_compress(pdf_path, output_path, compression_level):
                if os.path.getsize(output_path) < original_size:
                    print(f"[optimize_pdf] '{pdf_path}': Ghostscript succeeded and reduced size, done")
                    return True
                print(f"[optimize_pdf] '{pdf_path}': Ghostscript succeeded but didn't reduce size, trying PyMuPDF/pikepdf anyway")
            else:
                print(f"[optimize_pdf] '{pdf_path}': Ghostscript unavailable or failed, falling back to PyMuPDF")

            import fitz

            images_touched = False
            pdf_document = fitz.open(pdf_path)

            if compress_images:
                dpi_map = {"high": 150, "normal": 120, "very_reduced": 96}
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                target_dpi = dpi_map.get(compression_level, 120)
                jpeg_quality = quality_map.get(compression_level, 75)
                images_touched = self._downsample_pdf_images(pdf_document, target_dpi, jpeg_quality)

            try:
                pdf_document.subset_fonts()
            except Exception as e:
                print(f"[optimize_pdf] '{pdf_path}': subset_fonts skipped: {e}")

            if remove_metadata:
                pdf_document.set_metadata({})

            if images_touched:
                safe_options = {'garbage': 1, 'deflate': True, 'clean': False, 'deflate_fonts': True}
            else:
                safe_options = {'garbage': 4, 'deflate': True, 'clean': True, 'deflate_fonts': True}

            pymupdf_output = output_path + ".pymupdf_tmp.pdf"
            pdf_document.save(pymupdf_output, **safe_options)
            pdf_document.close()

            if self._validate_pdf(pymupdf_output):
                if not os.path.exists(output_path) or os.path.getsize(pymupdf_output) < os.path.getsize(output_path):
                    import shutil
                    shutil.move(pymupdf_output, output_path)
                    print(f"[optimize_pdf] '{pdf_path}': PyMuPDF save is smaller, using it")
                else:
                    os.remove(pymupdf_output)
                    print(f"[optimize_pdf] '{pdf_path}': PyMuPDF save didn't beat current best, discarding")
            else:
                print(f"[optimize_pdf] '{pdf_path}': PyMuPDF save failed validation, trying pikepdf")
                if os.path.exists(pymupdf_output):
                    os.remove(pymupdf_output)
                if not os.path.exists(output_path):
                    if not self._pikepdf_compress(pdf_path, output_path, compression_level):
                        print(f"[optimize_pdf] '{pdf_path}': pikepdf also failed, zero-touch resave")
                        zero_doc = fitz.open(pdf_path)
                        zero_doc.save(output_path, garbage=0, deflate=False, clean=False, incremental=False)
                        zero_doc.close()
                        if not self._validate_pdf(output_path):
                            print(f"[optimize_pdf] '{pdf_path}': source does not survive any re-save — copying original untouched.")
                            import shutil
                            shutil.copy2(pdf_path, output_path)
                            return False

            # Final guard: never hand back a file bigger than the input.
            if os.path.exists(output_path) and os.path.getsize(output_path) >= original_size:
                print(f"[optimize_pdf] '{pdf_path}': no method beat the original size, keeping original bytes")
                import shutil
                shutil.copy2(pdf_path, output_path)

            return True
        except Exception as e:
            print(f"PDF optimization error {pdf_path}: {e}")
            return False

    def _pikepdf_compress(self, pdf_path, output_path, compression_level):
        """Offline, pip-only fallback (bundles QPDF statically — no external
        binary, no download, no PATH/admin concerns). QPDF's repair is
        generally more tolerant of damaged xref tables than MuPDF's, and
        object-stream compaction + stream recompression gets real reduction
        on structurally messy 'normal' PDFs that PyMuPDF's save() bounced on.
        """
        try:
            import pikepdf
            with pikepdf.open(pdf_path) as pdf:
                pdf.remove_unreferenced_resources()
                pdf.save(
                    output_path,
                    compress_streams=True,
                    object_stream_mode=pikepdf.ObjectStreamMode.generate,
                    # "high" quality = leave existing stream compression alone;
                    # normal/very_reduced re-run flate compression, which can
                    # shrink streams that were saved uncompressed upstream.
                    recompress_flate=(compression_level != "high"),
                    linearize=False,
                )
            return self._validate_pdf(output_path)
        except ImportError:
            print("[pikepdf_compress] pikepdf not installed — run: pip install pikepdf")
            return False
        except Exception as e:
            print(f"[pikepdf_compress] {pdf_path}: {e}")
            return False

    def _ghostscript_compress(self, pdf_path, output_path, compression_level):
        """Ghostscript's pdfwrite device rebuilds the PDF from rendered content
        instead of patching the existing xref table — recovers files MuPDF's
        repair gives up on, and matches what most online PDF compressors use
        under the hood. Called by absolute path only, resolved by
        _get_ghostscript_bin() (detection only — see that method). Returns
        False if Ghostscript isn't installed; caller falls back to PyMuPDF.
        """
        gs_bin = self._get_ghostscript_bin()
        if not gs_bin:
            print(f"[ghostscript_compress] '{pdf_path}': no Ghostscript binary found "
                  f"(not on PATH, app-local folder, or common install paths)")
            return False
        print(f"[ghostscript_compress] '{pdf_path}': using binary at {gs_bin}")
        try:
            import subprocess
            settings_map = {"high": "/printer", "normal": "/ebook", "very_reduced": "/screen"}
            pdf_settings = settings_map.get(compression_level, "/ebook")
            tmp_out = output_path + ".gs_tmp.pdf"
            cmd = [gs_bin, "-o", tmp_out, "-sDEVICE=pdfwrite",
                   f"-dPDFSETTINGS={pdf_settings}", "-dNOPAUSE", "-dBATCH", "-dQUIET", pdf_path]
            # On Windows, subprocess would otherwise briefly flash a console
            # window for each conversion since this app has no console of its
            # own; CREATE_NO_WINDOW suppresses that. The attribute doesn't
            # exist on non-Windows platforms, hence the hasattr guard.
            _no_window = subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0
            result = subprocess.run(cmd, capture_output=True, timeout=180, creationflags=_no_window)
            if result.returncode != 0 or not os.path.exists(tmp_out):
                stderr_tail = result.stderr.decode(errors="replace")[-500:] if result.stderr else "(no stderr)"
                print(f"[ghostscript_compress] '{pdf_path}': gs exited {result.returncode}, tmp file exists={os.path.exists(tmp_out)}\n"
                      f"    stderr: {stderr_tail}")
                return False
            if not self._validate_pdf(tmp_out):
                print(f"[ghostscript_compress] '{pdf_path}': gs output failed PDF validation, discarding")
                os.remove(tmp_out)
                return False
            gs_size = os.path.getsize(tmp_out)
            orig_size = os.path.getsize(pdf_path)
            print(f"[ghostscript_compress] '{pdf_path}': gs succeeded, output size {gs_size} vs original {orig_size} "
                  f"({'smaller' if gs_size < orig_size else 'NOT smaller'})")
            import shutil
            shutil.move(tmp_out, output_path)
            return True
        except subprocess.TimeoutExpired:
            print(f"[ghostscript_compress] '{pdf_path}': gs timed out after 180s")
            return False
        except Exception as e:
            print(f"[ghostscript_compress] {pdf_path}: {e}")
            return False

    def _get_ghostscript_bin(self):
        """Locate a usable Ghostscript binary, in order:
        1. Already resolved this session (cached on self).
        2. A previously-installed app-local copy (self._gs_app_dir()).
        3. On the system PATH or common install locations.
        Detection only — no download, no auto-install. Returns an absolute
        path, or None if Ghostscript isn't installed anywhere findable, in
        which case the caller falls back to PyMuPDF/pikepdf.
        """
        if getattr(self, "_gs_bin_cache", None):
            return self._gs_bin_cache

        import shutil as _sh, platform, glob

        app_dir = self._gs_app_dir()
        local_candidates = {
            "Windows": [os.path.join(app_dir, "bin", "gswin64c.exe"),
                        os.path.join(app_dir, "bin", "gswin32c.exe")],
            "Darwin":  [os.path.join(app_dir, "bin", "gs")],
            "Linux":   [os.path.join(app_dir, "bin", "gs")],
        }.get(platform.system(), [])

        for c in local_candidates:
            if os.path.isfile(c):
                self._gs_bin_cache = c
                return c

        for name in ("gs", "gswin64c", "gswin32c"):
            found = _sh.which(name)
            if found:
                self._gs_bin_cache = found
                return found

        system_candidates = [
            r"C:\Program Files\gs\gs*\bin\gswin64c.exe",
            r"C:\Program Files (x86)\gs\gs*\bin\gswin32c.exe",
            "/usr/bin/gs", "/usr/local/bin/gs", "/opt/homebrew/bin/gs",
        ]
        for pattern in system_candidates:
            for match in glob.glob(pattern):
                if os.path.isfile(match):
                    self._gs_bin_cache = match
                    return match

        return None

    def check_ghostscript_status(self):
        """Diagnostic helper — call this on its own (e.g. from a settings/about
        screen) to know exactly what this app can and can't do, without
        running an actual PDF through it. Returns a dict rather than printing,
        so the caller can show it in the UI. Detection only, no download.
        """
        import shutil as _sh, subprocess, platform

        status = {"found": False, "path": None, "version": None,
                  "source": None, "error": None}

        cached = getattr(self, "_gs_bin_cache", None)
        if cached and os.path.isfile(cached):
            status["source"] = "cached"
            candidate = cached
        else:
            candidate = _sh.which("gs") or _sh.which("gswin64c") or _sh.which("gswin32c")
            if candidate:
                status["source"] = "system PATH"
            else:
                app_dir = self._gs_app_dir()
                for c in (os.path.join(app_dir, "bin", "gswin64c.exe"),
                          os.path.join(app_dir, "bin", "gswin32c.exe"),
                          os.path.join(app_dir, "bin", "gs")):
                    if os.path.isfile(c):
                        candidate = c
                        status["source"] = "app-local"
                        break

        if not candidate:
            status["error"] = "Not found on PATH, in app-local folder, or common install paths."
            return status

        # Presence on disk isn't proof it runs — confirm with a real call.
        try:
            _no_window = subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0
            result = subprocess.run([candidate, "--version"], capture_output=True,
                                     timeout=10, text=True, creationflags=_no_window)
            if result.returncode == 0:
                status["found"] = True
                status["path"] = candidate
                status["version"] = result.stdout.strip()
            else:
                status["error"] = f"Found at {candidate} but exited with code {result.returncode}: {result.stderr.strip()}"
        except Exception as e:
            status["error"] = f"Found at {candidate} but failed to execute: {e}"

        return status

    def _gs_app_dir(self):
        """App-local folder Ghostscript may already be installed to (if the
        user installed it there manually, or a previous version of this app
        did). Only ever checked now, never written to automatically."""
        base = getattr(self, "app_data_dir", None) or os.path.join(
            os.environ.get("LOCALAPPDATA") or os.path.expanduser("~"), ".file_converter_app")
        return os.path.join(base, "ghostscript")

    def _validate_pdf(self, path):
        """Reopen a saved PDF and confirm it's actually readable before trusting it.

        Checks every page, not just the first — corruption from xref surgery
        tends to cluster in high object numbers (later pages, embedded fonts,
        images), so a page-0-only check gives false positives on exactly the
        documents most at risk here.
        """
        try:
            import fitz
            check = fitz.open(path)
            if check.page_count == 0:
                check.close()
                return False
            for i in range(check.page_count):
                try:
                    check.load_page(i).get_text()
                except Exception as e:
                    print(f"PDF validation failed on page {i} of {path}: {e}")
                    check.close()
                    return False
            check.close()
            return True
        except Exception as e:
            print(f"PDF validation failed for {path}: {e}")
            return False

    def _downsample_pdf_images(self, pdf_document, target_dpi, jpeg_quality):
        """Resample embedded raster images to target_dpi and re-encode as JPEG.

        Effective DPI is based on the image's rendered size on the page, not its
        raw pixel dimensions. Images already at or below target_dpi are left
        unchanged. Vector content, masks (SMask/stencil), and CMYK images are
        skipped to avoid rendering issues.
        """

        from PIL import Image as PILImage
        import io as _io

        seen_xrefs = set()
        any_replaced = False

        for page in pdf_document:
            try:
                img_list = page.get_images(full=True)
            except Exception:
                continue

            for img in img_list:
                xref = img[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)

                # Skip images used as soft/stencil masks elsewhere.
                # img[1] is this image's SMask xref, but this xref may itself be used as a mask.
                try:
                    if pdf_document.xref_get_key(xref, "ImageMask")[1] == "true":
                        continue
                except Exception:
                    pass

                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                rect = rects[0]
                disp_w_in = rect.width / 72.0
                disp_h_in = rect.height / 72.0
                if disp_w_in <= 0 or disp_h_in <= 0:
                    continue

                try:
                    base_image = pdf_document.extract_image(xref)
                except Exception:
                    continue
                if not base_image:
                    continue

                px_w = base_image.get("width")
                px_h = base_image.get("height")
                img_bytes = base_image.get("image")
                if not px_w or not px_h or not img_bytes:
                    continue
                if base_image.get("smask"):
                    # Preserve alpha/soft mask.
                    continue

                if base_image.get("colorspace") == 4:  # CMYK
                    # Skip CMYK images.
                    continue

                current_dpi = px_w / disp_w_in

                # Some PDFs embed images at 1 px = 1 pt, making huge images appear ~72 DPI.
                # DPI-based checks alone would skip them, so also enforce an absolute pixel limit.
                abs_max_px = target_dpi * 20  # Allows up to a 20" side at target DPI
                dpi_over_target = current_dpi > target_dpi * 1.1
                pixel_dims_excessive = max(px_w, px_h) > abs_max_px

                if not dpi_over_target and not pixel_dims_excessive:
                    continue

                if dpi_over_target:
                    scale = target_dpi / current_dpi
                else:
                    scale = abs_max_px / max(px_w, px_h)

                new_w = max(1, int(px_w * scale))
                new_h = max(1, int(px_h * scale))

                try:
                    with PILImage.open(_io.BytesIO(img_bytes)) as pil_img:
                        if pil_img.mode in ('RGBA', 'P', 'LA'):
                            pil_img = pil_img.convert('RGB')
                        resized = pil_img.resize((new_w, new_h), PILImage.LANCZOS)
                        buf = _io.BytesIO()
                        resized.save(buf, format='JPEG', quality=jpeg_quality, optimize=True)
                        new_bytes = buf.getvalue()

                    if len(new_bytes) < len(img_bytes):
                        page.replace_image(xref, stream=new_bytes)
                        any_replaced = True
                except Exception as e:
                    print(f"[pdf downsample] xref {xref} skipped: {e}")
                    continue

        return any_replaced

    def optimize_word_file(self, word_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            from docx import Document
            import io as _io
            doc = Document(word_path)
            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage
                    for rel in doc.part.rels.values():
                        # reltype is a full URI (e.g. .../relationships/image);
                        # substring match is intentional, not a typo.
                        if "image" in rel.reltype:
                            try:
                                blob = rel.target_part.blob
                                ext  = Path(rel.target_part.partname).suffix.lower()
                                if ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff'):
                                    buf_in  = _io.BytesIO(blob)
                                    buf_out = _io.BytesIO()
                                    with PILImage.open(buf_in) as img:
                                        img.convert('RGB').save(buf_out, format='JPEG', quality=jpeg_quality, optimize=True)
                                    new_blob = buf_out.getvalue()
                                    if len(new_blob) < len(blob):
                                        # python-docx has no public setter for a part's
                                        # blob; writing the private attribute directly
                                        # is the documented workaround for in-place
                                        # image replacement.
                                        rel.target_part._blob = new_blob
                            except Exception:
                                pass
                except ImportError:
                    pass
            if remove_metadata:
                doc.core_properties.title = ""
                doc.core_properties.author = ""
                doc.core_properties.subject = ""
                doc.core_properties.keywords = ""
                doc.core_properties.comments = ""
                doc.core_properties.last_modified_by = ""
            doc.save(output_path)
            return True
        except Exception as e:
            print(f"Word optimization error {word_path}: {e}")
            return False

    def optimize_powerpoint_file(self, ppt_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(ppt_path)
            slides_to_remove = [i for i, slide in enumerate(prs.slides) if not slide.shapes]
            # python-pptx has no public API to delete a slide; removing its
            # entry from the slide ID list (the documented workaround) is
            # what actually drops it from the deck.
            xml_slides = prs.slides._sldIdLst
            for i in reversed(slides_to_remove):
                xml_slides.remove(xml_slides[i])
            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            # 13 == MSO_SHAPE_TYPE.PICTURE. Using the raw int
                            # avoids importing the enum for one comparison.
                            if shape.shape_type == 13:
                                try:
                                    img_part = shape.image
                                    blob = img_part.blob
                                    ext  = img_part.ext.lower()
                                    if ext in ("jpg", "jpeg", "png", "bmp", "tiff"):
                                        buf_in  = io.BytesIO(blob)
                                        buf_out = io.BytesIO()
                                        with PILImage.open(buf_in) as img:
                                            rgb = img.convert("RGB")
                                            rgb.save(buf_out, format="JPEG", quality=jpeg_quality, optimize=True)
                                        new_blob = buf_out.getvalue()
                                        if len(new_blob) < len(blob):
                                            img_part._blob = new_blob
                                except Exception:
                                    pass
                except ImportError:
                    pass
            if remove_metadata:
                prs.core_properties.title = ""
                prs.core_properties.author = ""
                prs.core_properties.subject = ""
                prs.core_properties.keywords = ""
                prs.core_properties.comments = ""
                prs.core_properties.last_modified_by = ""
            prs.save(output_path)
            return True
        except Exception as e:
            print(f"PowerPoint optimization error {ppt_path}: {e}")
            return False

    def optimize_av_file(self, src_path, output_path, quality_level, media_type):
        try:
            import shutil, sys, subprocess
            ffmpeg_bin = shutil.which("ffmpeg")
            if not ffmpeg_bin:
                candidates = [
                    r"C:\ffmpeg\bin\ffmpeg.exe", r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
                    r"C:\Program Files (x86)\ffmpeg\bin\ffmpeg.exe",
                    os.path.join(os.environ.get("LOCALAPPDATA", ""), "ffmpeg", "bin", "ffmpeg.exe"),
                    os.path.join(os.environ.get("APPDATA", ""), "ffmpeg", "bin", "ffmpeg.exe"),
                    os.path.join(getattr(sys, "_MEIPASS", ""), "ffmpeg.exe"),
                    "/usr/bin/ffmpeg", "/usr/local/bin/ffmpeg",
                ]
                for c in candidates:
                    if c and os.path.isfile(c):
                        ffmpeg_bin = c
                        break
            if not ffmpeg_bin:
                return False
            ext = Path(output_path).suffix.lower().lstrip(".")
            if media_type == 'audio':
                AUDIO_PRESETS = {
                    "mp3": {0: ["-codec:a", "libmp3lame", "-q:a", "2", "-ar", "44100"], 1: ["-codec:a", "libmp3lame", "-q:a", "4", "-ar", "44100"], 2: ["-codec:a", "libmp3lame", "-q:a", "7", "-ar", "44100"]},
                    "aac": {0: ["-codec:a", "aac", "-b:a", "192k", "-ar", "44100"], 1: ["-codec:a", "aac", "-b:a", "128k", "-ar", "44100"], 2: ["-codec:a", "aac", "-b:a", "96k", "-ar", "44100"]},
                    "ogg": {0: ["-codec:a", "libvorbis", "-q:a", "6"], 1: ["-codec:a", "libvorbis", "-q:a", "4"], 2: ["-codec:a", "libvorbis", "-q:a", "2"]},
                    "flac": {0: ["-codec:a", "flac", "-compression_level", "5"], 1: ["-codec:a", "flac", "-compression_level", "8"], 2: ["-codec:a", "flac", "-compression_level", "12"]},
                    "wav": {0: ["-codec:a", "pcm_s16le", "-ar", "44100"], 1: ["-codec:a", "pcm_s16le", "-ar", "44100"], 2: ["-codec:a", "pcm_s16le", "-ar", "22050"]},
                }
                default_audio = {0: ["-codec:a", "libmp3lame", "-q:a", "2"], 1: ["-codec:a", "libmp3lame", "-q:a", "4"], 2: ["-codec:a", "libmp3lame", "-q:a", "7"]}
                presets = AUDIO_PRESETS.get(ext, default_audio)
                args = presets.get(quality_level, presets[1])
            else:
                VIDEO_PRESETS = {
                    "mp4": {0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow", "-codec:a", "aac", "-b:a", "192k", "-movflags", "+faststart"], 1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium", "-codec:a", "aac", "-b:a", "128k", "-movflags", "+faststart"], 2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast", "-codec:a", "aac", "-b:a", "96k", "-movflags", "+faststart"]},
                    "mkv": {0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow", "-codec:a", "aac", "-b:a", "192k"], 1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium", "-codec:a", "aac", "-b:a", "128k"], 2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast", "-codec:a", "aac", "-b:a", "96k"]},
                    "webm": {0: ["-codec:v", "libvpx-vp9", "-crf", "24", "-b:v", "0", "-codec:a", "libopus", "-b:a", "160k"], 1: ["-codec:v", "libvpx-vp9", "-crf", "33", "-b:v", "0", "-codec:a", "libopus", "-b:a", "128k"], 2: ["-codec:v", "libvpx-vp9", "-crf", "42", "-b:v", "0", "-codec:a", "libopus", "-b:a", "96k"]},
                }
                default_video = {0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow", "-codec:a", "aac", "-b:a", "192k"], 1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium", "-codec:a", "aac", "-b:a", "128k"], 2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast", "-codec:a", "aac", "-b:a", "96k"]}
                presets = VIDEO_PRESETS.get(ext, default_video)
                args = presets.get(quality_level, presets[1])
            cmd = [ffmpeg_bin, "-y", "-i", src_path] + args + [output_path]
            _no_window = subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0
            result = subprocess.run(cmd, capture_output=True, timeout=3600, creationflags=_no_window)
            if result.returncode != 0:
                return False
            if os.path.exists(output_path):
                orig_size = os.path.getsize(src_path)
                new_size  = os.path.getsize(output_path)
                # 5% tolerance rather than a strict >=: re-encoding can add a
                # few % of container/codec overhead on already-compressed
                # source files without that being a real regression worth
                # discarding the transcode over.
                if new_size >= orig_size * 1.05:
                    shutil.copy2(src_path, output_path)
            return True
        except Exception as e:
            print(f"[optimize_av] Error {src_path}: {e}")
            return False

    def optimize_web_file(self, src_path, output_path, file_ext):
        try:
            with open(src_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            if file_ext == '.json':
                import json as _json
                data = _json.loads(content)
                minified = _json.dumps(data, ensure_ascii=False, separators=(',', ':'))
            else:
                import re as _re
                minified = _re.sub(r'<!--.*?-->', '', content, flags=_re.DOTALL)
                minified = _re.sub(r'>\s+<', '><', minified)
                minified = _re.sub(r'[ \t]{2,}', ' ', minified)
                lines = [l.strip() for l in minified.splitlines()]
                minified = '\n'.join(l for l in lines if l)
            orig_bytes = content.encode('utf-8')
            new_bytes  = minified.encode('utf-8')
            if len(new_bytes) < len(orig_bytes):
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(minified)
            else:
                import shutil as _sh
                _sh.copy2(src_path, output_path)
            return True
        except Exception as e:
            print(f"Web file optimization error {src_path}: {e}")
            return False

    def optimize_epub_file(self, src_path, output_path, compress_images, quality_level):
        try:
            import zipfile, io as _io
            quality_map = {0: 85, 1: 75, 2: 55}
            img_quality = quality_map.get(quality_level, 75)
            IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp'}
            buf_out = _io.BytesIO()
            with zipfile.ZipFile(src_path, 'r') as zin, zipfile.ZipFile(buf_out, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    ext  = Path(item.filename).suffix.lower()
                    if compress_images and ext in IMAGE_EXTS and len(data) > 2048:
                        try:
                            from PIL import Image as PILImage
                            buf_img = _io.BytesIO(data)
                            buf_new = _io.BytesIO()
                            with PILImage.open(buf_img) as img:
                                if ext in ('.jpg', '.jpeg'):
                                    img.convert('RGB').save(buf_new, format='JPEG', quality=img_quality, optimize=True)
                                elif ext == '.png':
                                    compress_lvl = max(1, min(9, int(quality_level * 3) + 1))
                                    img.save(buf_new, format='PNG', optimize=True, compress_level=compress_lvl)
                                else:
                                    img.convert('RGB').save(buf_new, format='JPEG', quality=img_quality, optimize=True)
                            new_data = buf_new.getvalue()
                            if len(new_data) < len(data):
                                data = new_data
                        except Exception:
                            pass
                    if item.filename == 'mimetype':
                        zout.writestr(item, data, compress_type=zipfile.ZIP_STORED)
                    else:
                        zout.writestr(item, data)
            new_bytes  = buf_out.getvalue()
            orig_bytes = os.path.getsize(src_path)
            if len(new_bytes) < orig_bytes:
                with open(output_path, 'wb') as f:
                    f.write(new_bytes)
            else:
                import shutil as _sh
                _sh.copy2(src_path, output_path)
            return True
        except Exception as e:
            print(f"EPUB optimization error {src_path}: {e}")
            return False

    def optimize_excel_file(self, xlsx_path, output_path, compression_level, remove_metadata):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(xlsx_path)
            if remove_metadata:
                wb.properties.title = ""
                wb.properties.creator = ""
                wb.properties.subject = ""
                wb.properties.keywords = ""
                wb.properties.description = ""
                wb.properties.lastModifiedBy = ""
            # A freshly-created/never-used worksheet reports max_row=max_column=1
            # with an empty A1 — that's the signature this checks for. A sheet
            # that's genuinely just one populated cell is indistinguishable
            # from this and won't be flagged, since its A1 value isn't None.
            empty_sheets = [ws.title for ws in wb.worksheets if ws.max_row == 1 and ws.max_column == 1 and ws.cell(1, 1).value is None]
            for name in empty_sheets:
                # A workbook can't have zero sheets — guard against deleting
                # the last one even if it's "empty".
                if len(wb.sheetnames) > 1:
                    del wb[name]
            wb.save(output_path)
            return True
        except Exception as e:
            print(f"Excel optimization error {xlsx_path}: {e}")
            return False

    def optimize_image_file(self, img_path, output_path, quality_level):
        try:
            from PIL import Image as PILImage
            quality_map = {0: 85, 1: 75, 2: 55}
            quality = quality_map.get(quality_level, 75)
            ext_map = {".jpg": "JPEG", ".jpeg": "JPEG", ".png": "PNG", ".bmp": "PNG", ".tiff": "TIFF", ".webp": "WEBP", ".gif": "GIF"}
            src_ext  = Path(img_path).suffix.lower()
            dst_ext  = Path(output_path).suffix.lower()
            fmt_out  = ext_map.get(dst_ext, ext_map.get(src_ext, "JPEG"))
            with PILImage.open(img_path) as img:
                exif_data = None
                try:
                    exif_data = img.info.get("exif")
                except Exception:
                    pass
                if fmt_out == "JPEG":
                    rgb = img.convert("RGB")
                    save_kwargs = {"quality": quality, "optimize": True}
                    if exif_data:
                        save_kwargs["exif"] = exif_data
                    rgb.save(output_path, format="JPEG", **save_kwargs)
                elif fmt_out == "PNG":
                    # PNG has no "quality" concept — quality_level (0/1/2) is
                    # remapped onto PIL's compress_level range (0-9, lossless
                    # either way; higher just spends more CPU for a smaller
                    # file). 0->0(min,clamped to 1), 1->3, 2->6.
                    compress = max(1, min(9, int(quality_level * 3)))
                    img.save(output_path, format="PNG", optimize=True, compress_level=compress)
                elif fmt_out == "WEBP":
                    # method=6 = slowest/best compression PIL's WEBP encoder
                    # offers; fine here since this runs once per file, not
                    # in a hot loop.
                    img.save(output_path, format="WEBP", quality=quality, method=6)
                else:
                    img.save(output_path, format=fmt_out)
            return True
        except Exception as e:
            print(f"Image optimization error {img_path}: {e}")
            return False