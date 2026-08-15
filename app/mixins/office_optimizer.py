"""Office document optimization (Word / PowerPoint / Excel)."""

from __future__ import annotations

from pathlib import Path


class OfficeDocOptimizerMixin:
    """Word / PowerPoint / Excel optimization methods."""

    def optimize_word_file(self, word_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            import io as _io

            from docx import Document

            doc = Document(word_path)
            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage

                    for rel in doc.part.rels.values():
                        # reltype is a full URI (e.g. .../relationships/image);
                        # substring match is intentional, not a typo.
                        if "image" in rel.reltype:
                            try:
                                blob = rel.target_part.blob
                                ext = Path(rel.target_part.partname).suffix.lower()
                                if ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff"):
                                    buf_in = _io.BytesIO(blob)
                                    buf_out = _io.BytesIO()
                                    with PILImage.open(buf_in) as img:
                                        img.convert("RGB").save(
                                            buf_out, format="JPEG", quality=jpeg_quality, optimize=True
                                        )
                                    new_blob = buf_out.getvalue()
                                    if len(new_blob) < len(blob):
                                        # python-docx has no public setter for a part's
                                        # blob; writing the private attribute directly
                                        # is the documented workaround for in-place
                                        # image replacement.
                                        rel.target_part._blob = new_blob
                            except Exception:
                                pass
                except ImportError:
                    pass
            if remove_metadata:
                doc.core_properties.title = ""
                doc.core_properties.author = ""
                doc.core_properties.subject = ""
                doc.core_properties.keywords = ""
                doc.core_properties.comments = ""
                doc.core_properties.last_modified_by = ""
            doc.save(output_path)
            return True
        except Exception as e:
            print(f"Word optimization error {word_path}: {e}")
            return False

    def optimize_powerpoint_file(self, ppt_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            import io

            from pptx import Presentation

            prs = Presentation(ppt_path)
            slides_to_remove = [i for i, slide in enumerate(prs.slides) if not slide.shapes]
            # python-pptx has no public API to delete a slide; removing its
            # entry from the slide ID list (the documented workaround) is
            # what actually drops it from the deck.
            xml_slides = prs.slides._sldIdLst
            for i in reversed(slides_to_remove):
                xml_slides.remove(xml_slides[i])
            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage

                    for slide in prs.slides:
                        for shape in slide.shapes:
                            # 13 == MSO_SHAPE_TYPE.PICTURE. Using the raw int
                            # avoids importing the enum for one comparison.
                            if shape.shape_type == 13:
                                try:
                                    img_part = shape.image
                                    blob = img_part.blob
                                    ext = img_part.ext.lower()
                                    if ext in ("jpg", "jpeg", "png", "bmp", "tiff"):
                                        buf_in = io.BytesIO(blob)
                                        buf_out = io.BytesIO()
                                        with PILImage.open(buf_in) as img:
                                            rgb = img.convert("RGB")
                                            rgb.save(buf_out, format="JPEG", quality=jpeg_quality, optimize=True)
                                        new_blob = buf_out.getvalue()
                                        if len(new_blob) < len(blob):
                                            img_part._blob = new_blob
                                except Exception:
                                    pass
                except ImportError:
                    pass
            if remove_metadata:
                prs.core_properties.title = ""
                prs.core_properties.author = ""
                prs.core_properties.subject = ""
                prs.core_properties.keywords = ""
                prs.core_properties.comments = ""
                prs.core_properties.last_modified_by = ""
            prs.save(output_path)
            return True
        except Exception as e:
            print(f"PowerPoint optimization error {ppt_path}: {e}")
            return False

    def optimize_excel_file(self, xlsx_path, output_path, compression_level, remove_metadata):
        try:
            import openpyxl

            wb = openpyxl.load_workbook(xlsx_path)
            if remove_metadata:
                wb.properties.title = ""
                wb.properties.creator = ""
                wb.properties.subject = ""
                wb.properties.keywords = ""
                wb.properties.description = ""
                wb.properties.lastModifiedBy = ""
            # A freshly-created/never-used worksheet reports max_row=max_column=1
            # with an empty A1 — that's the signature this checks for. A sheet
            # that's genuinely just one populated cell is indistinguishable
            # from this and won't be flagged, since its A1 value isn't None.
            empty_sheets = [
                ws.title
                for ws in wb.worksheets
                if ws.max_row == 1 and ws.max_column == 1 and ws.cell(1, 1).value is None
            ]
            for name in empty_sheets:
                # A workbook can't have zero sheets — guard against deleting
                # the last one even if it's "empty".
                if len(wb.sheetnames) > 1:
                    del wb[name]
            wb.save(output_path)
            return True
        except Exception as e:
            print(f"Excel optimization error {xlsx_path}: {e}")
            return False
