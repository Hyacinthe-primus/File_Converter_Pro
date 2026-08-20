"""Native PPTX -> PDF rendering (python-pptx + reportlab)."""

from __future__ import annotations

import os
import tempfile

from converter.helpers import _safe_html


class PptxNativeConverters:
    """Native PowerPoint -> PDF rendering."""

    def _pptx_to_pdf_native(self, src, dst):
        """
        Native PPTX → PDF using python-pptx + reportlab + matplotlib.
        Handles: text, bullets, images, tables, charts (bar/line/pie/donut).
        All python-pptx API calls are wrapped in try/except.
        """
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER, TA_RIGHT
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import HRFlowable, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from reportlab.platypus import Image as RLImage

        prs = Presentation(src)
        try:
            sw = int(prs.slide_width) or 9144000
        except Exception:
            sw = 9144000
        try:
            sh = int(prs.slide_height) or 6858000
        except Exception:
            sh = 6858000

        ratio = float(sh) / float(sw)
        pw = landscape(A4)[0]
        ph = pw * ratio
        margin = 14 * mm
        inner_w = pw - 2 * margin

        S = getSampleStyleSheet()

        def sty(name, **kw):
            p = kw.pop("parent", S["Normal"])
            return ParagraphStyle(name, parent=p, **kw)

        ST = {
            "title": sty(
                "PT",
                parent=S["Heading1"],
                fontSize=20,
                leading=24,
                spaceAfter=8,
                textColor=colors.HexColor("#1a1a2e"),
                alignment=TA_CENTER,
            ),
            "h2": sty(
                "PH2", parent=S["Heading2"], fontSize=15, leading=19, spaceAfter=5, textColor=colors.HexColor("#1e3a5f")
            ),
            "body": sty("PB", fontSize=11, leading=16, spaceAfter=4),
            "bul0": sty("PB0", fontSize=11, leading=16, spaceAfter=3, leftIndent=12),
            "bul1": sty("PB1", fontSize=10, leading=15, spaceAfter=2, leftIndent=24, textColor=colors.HexColor("#333")),
            "bul2": sty("PB2", fontSize=9, leading=14, spaceAfter=2, leftIndent=36, textColor=colors.HexColor("#555")),
            "slide_n": sty("PSN", fontSize=8, leading=10, textColor=colors.HexColor("#aaa"), alignment=TA_RIGHT),
            "td": sty("PTd", fontSize=9, leading=13),
            "th": sty("PTh", fontSize=9, leading=13, fontName="Helvetica-Bold"),
            "caption": sty("PC", fontSize=8, leading=11, textColor=colors.HexColor("#777"), alignment=TA_CENTER),
        }

        BULLETS = set("•–-*◦▪▸")
        story = []
        tmp_imgs = []

        def _safe(txt):
            return _safe_html(str(txt or ""))

        def _add_img_blob(blob, ext):
            try:
                tf = tempfile.NamedTemporaryFile(suffix=f".{ext or 'png'}", delete=False)
                tf.write(blob)
                tf.close()
                tmp_imgs.append(tf.name)
                rl = RLImage(tf.name)
                scale = min(inner_w / rl.imageWidth, (ph - 2 * margin) * 0.6 / rl.imageHeight, 1.0)
                rl.drawWidth = rl.imageWidth * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign = "CENTER"
                story.append(Spacer(1, 4))
                story.append(rl)
                story.append(Spacer(1, 4))
            except Exception:
                pass

        def _render_chart(shape):
            """Extract chart data and render as matplotlib image."""
            try:
                import matplotlib

                matplotlib.use("Agg")
                import matplotlib.pyplot as plt
                import numpy as np
                from pptx.enum.chart import XL_CHART_TYPE

                chart = shape.chart
                ctype = chart.chart_type

                NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                cxml = chart._element

                series_list = []
                for ser in cxml.findall(f".//{{{NS_C}}}ser"):
                    name = ""
                    v_el = ser.find(f".//{{{NS_C}}}tx//{{{NS_C}}}v")
                    if v_el is not None:
                        name = v_el.text or ""

                    cats = [
                        el.find(f"{{{NS_C}}}v").text
                        for el in ser.findall(f".//{{{NS_C}}}cat//{{{NS_C}}}pt")
                        if el.find(f"{{{NS_C}}}v") is not None
                    ]

                    vals = []
                    for pt in ser.findall(f".//{{{NS_C}}}val//{{{NS_C}}}pt"):
                        v = pt.find(f"{{{NS_C}}}v")
                        if v is not None:
                            try:
                                vals.append(float(v.text))
                            except Exception:
                                vals.append(0.0)

                    pt_colors = {}
                    for dpt in ser.findall(f"{{{NS_C}}}dPt"):
                        idx_el = dpt.find(f"{{{NS_C}}}idx")
                        clr_el = dpt.find(f".//{{{NS_A}}}srgbClr")
                        if idx_el is not None and clr_el is not None:
                            idx = int(idx_el.get("val", 0))
                            pt_colors[idx] = "#" + clr_el.get("val", "4472C4")

                    ser_clr = None
                    clr_el = ser.find(f".//{{{NS_A}}}srgbClr")
                    if clr_el is not None:
                        ser_clr = "#" + clr_el.get("val", "4472C4")

                    series_list.append(
                        {
                            "name": name,
                            "cats": cats,
                            "vals": vals,
                            "pt_colors": pt_colors,
                            "ser_color": ser_clr,
                        }
                    )

                if not series_list:
                    return False

                fig_w = float(inner_w) / 72 / 1.333 * 1.5
                fig_h = fig_w * 0.6
                fig, ax = plt.subplots(figsize=(fig_w, fig_h))
                fig.patch.set_facecolor("white")
                ax.set_facecolor("#f8f9fa")

                is_pie = ctype in (
                    XL_CHART_TYPE.PIE,
                    XL_CHART_TYPE.PIE_EXPLODED,
                    XL_CHART_TYPE.DOUGHNUT,
                    XL_CHART_TYPE.DOUGHNUT_EXPLODED,
                )
                is_bar = ctype in (
                    XL_CHART_TYPE.BAR_CLUSTERED,
                    XL_CHART_TYPE.BAR_STACKED,
                    XL_CHART_TYPE.BAR_STACKED_100,
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    XL_CHART_TYPE.COLUMN_STACKED,
                )
                is_line = ctype in (
                    XL_CHART_TYPE.LINE,
                    XL_CHART_TYPE.LINE_MARKERS,
                )

                ser0 = series_list[0]

                if is_pie:
                    vals_ = ser0["vals"]
                    labels = ser0["cats"] or [f"Cat {i + 1}" for i in range(len(vals_))]
                    clrs = [ser0["pt_colors"].get(i, None) for i in range(len(vals_))]
                    # Fill None colors with defaults
                    default_clrs = plt.rcParams["axes.prop_cycle"].by_key()["color"]
                    clrs = [c if c else default_clrs[i % len(default_clrs)] for i, c in enumerate(clrs)]

                    wedge_kw = {"width": 0.55} if "DOUGHNUT" in str(ctype) else {}
                    wedges, texts, autotexts = ax.pie(
                        vals_,
                        labels=None,
                        colors=clrs,
                        autopct="%1.1f%%",
                        startangle=90,
                        wedgeprops=wedge_kw,
                        pctdistance=0.75,
                    )
                    for at in autotexts:
                        at.set_fontsize(8)
                    ax.legend(wedges, labels, loc="center left", bbox_to_anchor=(1, 0.5), fontsize=8)
                    if "DOUGHNUT" in str(ctype):
                        total = sum(vals_)
                        ax.text(0, 0, f"{total:.0f}", ha="center", va="center", fontsize=12, fontweight="bold")
                    ax.set_title(ser0["name"] or "Chart", fontsize=11, pad=10)

                elif is_bar or (not is_line):
                    cats = ser0["cats"] or [str(i + 1) for i in range(len(ser0["vals"]))]
                    x = np.arange(len(cats))
                    n_ser = len(series_list)
                    width = 0.8 / max(n_ser, 1)
                    offsets = np.linspace(-(n_ser - 1) * width / 2, (n_ser - 1) * width / 2, n_ser)
                    clr_cycle = plt.rcParams["axes.prop_cycle"].by_key()["color"]

                    for si2, ser2 in enumerate(series_list):
                        clr = ser2["ser_color"] or clr_cycle[si2 % len(clr_cycle)]
                        ax.bar(x + offsets[si2], ser2["vals"], width, label=ser2["name"], color=clr, alpha=0.88)

                    ax.set_xticks(x)
                    ax.set_xticklabels(cats, rotation=30, ha="right", fontsize=8)
                    ax.set_ylabel("Value", fontsize=9)
                    ax.tick_params(axis="y", labelsize=8)
                    ax.grid(axis="y", alpha=0.3, linestyle="--")
                    if any(s["name"] for s in series_list):
                        ax.legend(fontsize=8)

                else:
                    cats = ser0["cats"] or [str(i + 1) for i in range(len(ser0["vals"]))]
                    for si2, ser2 in enumerate(series_list):
                        ax.plot(cats, ser2["vals"], marker="o", label=ser2["name"], linewidth=1.8, markersize=4)
                    ax.set_xticks(range(len(cats)))
                    ax.set_xticklabels(cats, rotation=30, ha="right", fontsize=8)
                    ax.grid(alpha=0.3, linestyle="--")
                    ax.legend(fontsize=8)

                plt.tight_layout(pad=1.0)

                tf = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                fig.savefig(tf.name, dpi=150, bbox_inches="tight", facecolor="white")
                plt.close(fig)
                tf.close()
                tmp_imgs.append(tf.name)

                rl = RLImage(tf.name)
                scale = min(inner_w / rl.imageWidth, (ph - 2 * margin) * 0.55 / rl.imageHeight, 1.0)
                rl.drawWidth = rl.imageWidth * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign = "CENTER"
                story.append(Spacer(1, 4))
                story.append(rl)
                story.append(Spacer(1, 6))
                return True

            except Exception:
                try:
                    chart = shape.chart
                    story.append(Paragraph(f"[Chart — {chart.chart_type}]", ST["h2"]))
                    for ser in chart.series:
                        story.append(Paragraph(f"• {ser.name}", ST["body"]))
                except Exception:
                    story.append(Paragraph("[Chart]", ST["body"]))
                return False

        def _render_table(tbl_obj):
            try:
                rows_data = []
                for ri, row in enumerate(tbl_obj.rows):
                    cells = []
                    for ci, cell in enumerate(row.cells):
                        try:
                            txt = cell.text.strip()
                        except Exception:
                            txt = ""
                        ps = ST["th"] if ri == 0 else ST["td"]
                        cells.append(Paragraph(_safe(txt) or " ", ps))
                    rows_data.append(cells)
                if not rows_data:
                    return
                t = Table(rows_data, repeatRows=1)
                t.setStyle(
                    TableStyle(
                        [
                            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#ccc")),
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e8edf5")),
                            ("VALIGN", (0, 0), (-1, -1), "TOP"),
                            ("LEFTPADDING", (0, 0), (-1, -1), 4),
                            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                            ("TOPPADDING", (0, 0), (-1, -1), 3),
                            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                        ]
                    )
                )
                story.append(Spacer(1, 6))
                story.append(t)
                story.append(Spacer(1, 6))
            except Exception:
                pass

        def _render_text_frame(tf_obj, default_sty):
            for para in tf_obj.paragraphs:
                try:
                    text = para.text.strip()
                    if not text:
                        story.append(Spacer(1, 3))
                        continue
                    try:
                        level = int(para.level or 0)
                    except Exception:
                        level = 0
                    if level == 0:
                        is_bul = text[0] in BULLETS
                        ps = ST["bul0"] if is_bul else default_sty
                        prefix = "• " if is_bul and text[0] not in BULLETS else ""
                    elif level == 1:
                        ps, prefix = ST["bul1"], "  ◦ "
                    else:
                        ps, prefix = ST["bul2"], "    ▪ "
                    parts = []
                    for run in para.runs:
                        try:
                            rt = run.text
                            if not rt:
                                continue
                            s = _safe(rt)
                            try:
                                b, i = run.font.bold, run.font.italic
                                if b and i:
                                    s = f"<b><i>{s}</i></b>"
                                elif b:
                                    s = f"<b>{s}</b>"
                                elif i:
                                    s = f"<i>{s}</i>"
                            except Exception:
                                pass
                            parts.append(s)
                        except Exception:
                            pass
                    inner = prefix + ("".join(parts) if parts else _safe(prefix + text))
                    story.append(Paragraph(inner, ps))
                except Exception:
                    try:
                        raw = para.text.strip()
                        if raw:
                            story.append(Paragraph(_safe(raw), default_sty))
                    except Exception:
                        pass

        for sn, slide in enumerate(prs.slides, 1):
            shapes_sorted = sorted(
                slide.shapes,
                key=lambda sh: (
                    (sh.top or 0),
                    (sh.left or 0),
                ),
            )

            for shape in shapes_sorted:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            _add_img_blob(shape.image.blob, shape.image.ext)
                        except Exception:
                            pass
                        continue

                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        _render_chart(shape)
                        continue

                    if shape.has_table:
                        _render_table(shape.table)
                        continue

                    if not shape.has_text_frame:
                        continue
                    if not shape.text_frame.text.strip():
                        continue

                    is_title = shape.name.lower().startswith("title")

                    if is_title:
                        story.append(Paragraph(_safe(shape.text_frame.text.strip()), ST["title"]))
                        story.append(
                            HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#1e3a8a"), spaceAfter=6)
                        )
                    else:
                        _render_text_frame(shape.text_frame, ST["body"])

                except Exception:
                    pass

            story.append(Spacer(1, 6))
            story.append(Paragraph(f"— {sn} / {len(prs.slides)} —", ST["slide_n"]))
            if sn < len(prs.slides):
                story.append(PageBreak())

        if not story:
            story.append(Paragraph("(empty presentation)", ST["body"]))

        def _safe_title(p):
            try:
                return p.slides[0].shapes.title.text or ""
            except Exception:
                return ""

        SimpleDocTemplate(
            dst,
            pagesize=(pw, ph),
            leftMargin=margin,
            rightMargin=margin,
            topMargin=margin,
            bottomMargin=margin,
            title=_safe_title(prs),
        ).build(story)

        for f in tmp_imgs:
            try:
                os.remove(f)
            except Exception:
                pass
